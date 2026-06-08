#!/usr/bin/env python3
"""
Generate PowerPoint Presentation: Strengthening Mozambique's Cyber Resilience
Partnership between Bouquet Innovation and Infosec
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from datetime import datetime

# Color scheme
BOUQUET_BLUE = RGBColor(0, 61, 122)      # #003D7A
GOLD = RGBColor(212, 175, 55)             # #D4AF37
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK_GRAY = RGBColor(64, 64, 64)
LIGHT_GRAY = RGBColor(240, 240, 240)

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Set up blank slide layout
    blank_slide_layout = prs.slide_layouts[6]

    print("Building PowerPoint Presentation...")
    print("=" * 60)

    # Slide 1: Cover Slide
    print("Creating Slide 1: Cover Slide...")
    slide1 = add_cover_slide(prs, blank_slide_layout)

    # Slide 2: Executive Summary
    print("Creating Slide 2: Executive Summary...")
    slide2 = add_executive_summary(prs, blank_slide_layout)

    # Slide 3: About Bouquet Innovation
    print("Creating Slide 3: About Bouquet Innovation...")
    slide3 = add_bouquet_slide(prs, blank_slide_layout)

    # Slide 4: About Infosec
    print("Creating Slide 4: About Infosec...")
    slide4 = add_infosec_slide(prs, blank_slide_layout)

    # Slide 5: Cybersecurity Challenges
    print("Creating Slide 5: Cybersecurity Challenges...")
    slide5 = add_challenges_slide(prs, blank_slide_layout)

    # Slide 6: Why Awareness Matters
    print("Creating Slide 6: Why Awareness Matters...")
    slide6 = add_awareness_matters_slide(prs, blank_slide_layout)

    # Slide 7: Awareness Program Framework
    print("Creating Slide 7: Awareness Program Framework...")
    slide7 = add_awareness_framework_slide(prs, blank_slide_layout)

    # Slide 8: Cyber Intelligence Needs
    print("Creating Slide 8: Cyber Intelligence Needs...")
    slide8 = add_intelligence_needs_slide(prs, blank_slide_layout)

    # Slide 9: Cyber Intelligence Unit Model
    print("Creating Slide 9: Cyber Intelligence Unit Model...")
    slide9 = add_intelligence_unit_slide(prs, blank_slide_layout)

    # Slide 10: National Cyberspace Management
    print("Creating Slide 10: National Cyberspace Management...")
    slide10 = add_cyberspace_management_slide(prs, blank_slide_layout)

    # Slide 11: Technology & Tools
    print("Creating Slide 11: Technology & Tools...")
    slide11 = add_technology_slide(prs, blank_slide_layout)

    # Slide 12: Implementation Roadmap
    print("Creating Slide 12: Implementation Roadmap...")
    slide12 = add_roadmap_slide(prs, blank_slide_layout)

    # Slide 13: Expected Impact
    print("Creating Slide 13: Expected Impact...")
    slide13 = add_impact_slide(prs, blank_slide_layout)

    # Slide 14: Partnership Rationale
    print("Creating Slide 14: Partnership Rationale...")
    slide14 = add_partnership_rationale_slide(prs, blank_slide_layout)

    # Slide 15: Call to Action
    print("Creating Slide 15: Call to Action...")
    slide15 = add_call_to_action_slide(prs, blank_slide_layout)

    # Slide 16: Contact Information
    print("Creating Slide 16: Contact Information...")
    slide16 = add_contact_slide(prs, blank_slide_layout)

    # Save presentation
    filename = "Bouquet_Infosec_Mozambique_Partnership.pptx"
    prs.save(filename)
    print("=" * 60)
    print(f"✅ Presentation created successfully!")
    print(f"📁 File: {filename}")
    print(f"📊 Slides: 16")
    print(f"📅 Created: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    print("=" * 60)

def add_footer_and_formatting(slide, slide_number, total_slides=16):
    """Add footer and basic formatting to slide"""
    # Add footer text box
    footer_left = slide.shapes.add_textbox(Inches(0.3), Inches(7.0), Inches(9.4), Inches(0.4))
    footer_frame = footer_left.text_frame
    footer_frame.text = f"Slide {slide_number} of {total_slides} | Bouquet Innovation & Infosec | {datetime.now().year}"
    footer_frame.paragraphs[0].font.size = Pt(10)
    footer_frame.paragraphs[0].font.color.rgb = DARK_GRAY
    footer_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

def add_title_shape(slide, title_text, y_position=Inches(0.4)):
    """Add formatted title to slide"""
    title_box = slide.shapes.add_textbox(Inches(0.4), y_position, Inches(9.2), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_frame.word_wrap = True

    # Format title
    paragraph = title_frame.paragraphs[0]
    paragraph.font.size = Pt(48)
    paragraph.font.bold = True
    paragraph.font.color.rgb = BOUQUET_BLUE

    return title_box

def add_content_box(slide, content_list, x=Inches(0.4), y=Inches(1.3),
                    width=Inches(9.2), height=Inches(5.5)):
    """Add bullet points to slide"""
    text_box = slide.shapes.add_textbox(x, y, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(content_list):
        if i > 0:
            text_frame.add_paragraph()

        p = text_frame.paragraphs[i]
        p.text = item
        p.level = 0
        p.font.size = Pt(24)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
        p.space_after = Pt(6)

    return text_box

def add_shape_line(slide, x1, y1, x2, y2, color=BOUQUET_BLUE, width=Pt(3)):
    """Add a line to the slide"""
    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = width

def add_cover_slide(prs, blank_layout):
    """Slide 1: Cover Slide"""
    slide = prs.slides.add_slide(blank_layout)

    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BOUQUET_BLUE

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.text = "Strengthening Mozambique's Cyber Resilience"
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(1.2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_frame.text = "A Strategic Partnership Between Bouquet Innovation and Infosec"
    subtitle_frame.paragraphs[0].font.size = Pt(32)
    subtitle_frame.paragraphs[0].font.color.rgb = GOLD
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(0.8))
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = "Building National Cybersecurity Resilience for a Secure Digital Future"
    tagline_frame.paragraphs[0].font.size = Pt(20)
    tagline_frame.paragraphs[0].font.italic = True
    tagline_frame.paragraphs[0].font.color.rgb = WHITE
    tagline_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.6))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Bouquet Innovation  |  Infosec  |  Republic of Mozambique"
    footer_frame.paragraphs[0].font.size = Pt(16)
    footer_frame.paragraphs[0].font.color.rgb = LIGHT_GRAY
    footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_text = notes_slide.notes_text_frame
    notes_text.text = """This presentation outlines a transformative partnership designed to elevate
Mozambique's cybersecurity posture across three critical domains: awareness,
intelligence, and national cyberspace management. The partnership leverages
Bouquet Innovation's expertise in digital transformation and strategic innovation
with Infosec's world-class cybersecurity capabilities."""

    add_footer_and_formatting(slide, 1)
    return slide

def add_executive_summary(prs, blank_layout):
    """Slide 2: Executive Summary"""
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Add accent bar
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BOUQUET_BLUE
    accent.line.color.rgb = BOUQUET_BLUE

    # Title
    add_title_shape(slide, "Executive Summary: Partnership Overview")

    # Content
    content = [
        "✓ Objective: Establish Mozambique as a cyber-resilient nation through integrated awareness, intelligence, and infrastructure programs",
        "✓ Scope: Three complementary pillars addressing national-level cybersecurity challenges",
        "✓ Timeline: 24-month implementation roadmap with clear milestones",
        "✓ Impact: Enhanced government security, increased investor confidence, reduced cybercriminal activity",
        "✓ Partners: Bouquet Innovation (strategic delivery) + Infosec (technical expertise)"
    ]

    add_content_box(slide, content)

    # Add diagram description box
    diagram_box = slide.shapes.add_shape(1, Inches(0.4), Inches(5.8), Inches(9.2), Inches(1.0))
    diagram_box.fill.solid()
    diagram_box.fill.fore_color.rgb = LIGHT_GRAY
    diagram_box.line.color.rgb = BOUQUET_BLUE

    diagram_text = diagram_box.text_frame
    diagram_text.text = "[DIAGRAM: Central 'Cyber Resilience Goal' hub with three branches: Awareness Programs | Cyber Intelligence Unit | National Cyberspace Management]"
    diagram_text.paragraphs[0].font.size = Pt(14)
    diagram_text.paragraphs[0].font.italic = True
    diagram_text.paragraphs[0].font.color.rgb = DARK_GRAY
    diagram_text.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text = notes_slide.notes_text_frame
    notes_text.text = """This partnership is designed to address Mozambique's most pressing cybersecurity
vulnerabilities. Rather than treating cybersecurity as an IT problem, we approach it
as a national security and economic resilience imperative. The three pillars are
interdependent and designed to create a comprehensive security ecosystem. The
24-month timeline is aggressive but achievable with proper commitment and resources."""

    add_footer_and_formatting(slide, 2)
    return slide

def add_bouquet_slide(prs, blank_layout):
    """Slide 3: About Bouquet Innovation"""
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Accent bar
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BOUQUET_BLUE
    accent.line.color.rgb = BOUQUET_BLUE

    # Title
    add_title_shape(slide, "Bouquet Innovation: Digital Transformation Partner")

    # Content
    content = [
        "Expertise: Government modernization | Strategic cybersecurity initiatives | Digital service delivery | Technology deployment",
        "Regional Focus: Africa, with emphasis on Portuguese-speaking nations",
        "Core Competency: Translating complex technical requirements into government-ready programs",
        "",
        "✓ Deep understanding of African government structures",
        "✓ Proven ability to deliver large-scale digital initiatives",
        "✓ Strong relationships with government and international bodies",
        "✓ Track record of successful partnerships with leading security firms"
    ]

    add_content_box(slide, content, y=Inches(1.2))

    # Contact info box
    contact_box = slide.shapes.add_shape(1, Inches(0.4), Inches(6.0), Inches(9.2), Inches(0.8))
    contact_box.fill.solid()
    contact_box.fill.fore_color.rgb = LIGHT_GRAY
    contact_box.line.color.rgb = GOLD
    contact_box.line.width = Pt(2)

    contact_text = contact_box.text_frame
    contact_text.text = "Email: support@bouquet-innovation.net  |  Phone: +258 878 275 656"
    contact_text.paragraphs[0].font.size = Pt(16)
    contact_text.paragraphs[0].font.bold = True
    contact_text.paragraphs[0].font.color.rgb = BOUQUET_BLUE
    contact_text.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text = notes_slide.notes_text_frame
    notes_text.text = """Bouquet Innovation serves as the strategic integrator and delivery partner. Our role
is to understand Mozambique's unique needs, coordinate across government stakeholders,
manage implementation timelines, and ensure that Infosec's world-class capabilities
are deployed effectively within the Mozambican context."""

    add_footer_and_formatting(slide, 3)
    return slide

def add_infosec_slide(prs, blank_layout):
    """Slide 4: About Infosec"""
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Accent bar
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BOUQUET_BLUE
    accent.line.color.rgb = BOUQUET_BLUE

    # Title
    add_title_shape(slide, "Infosec: Global Cybersecurity Excellence")

    # Content
    content = [
        "Global Leader in: Cyber threat intelligence | Digital forensics | Incident response | Advanced threat research",
        "",
        "Expertise Areas:",
        "  • Cyber intelligence and threat monitoring  • Incident response and recovery",
        "  • Digital forensics and malware analysis  • Counter-terrorism and fraud support",
        "",
        "✓ Trusted advisor to governments and national security agencies",
        "✓ Proven capability in complex, high-stakes environments",
        "✓ Advanced threat research and attribution capabilities",
        "✓ 24/7 global incident response network"
    ]

    add_content_box(slide, content, y=Inches(1.2))

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text = notes_slide.notes_text_frame
    notes_text.text = """Infosec brings world-class technical expertise and proven government-level
operational experience. Their capabilities are not theoretical—they are battle-tested
in complex, real-world scenarios. This partnership allows Mozambique to access
enterprise-grade cybersecurity without the years of capability development that
would otherwise be required."""

    add_footer_and_formatting(slide, 4)
    return slide

def add_challenges_slide(prs, blank_layout):
    """Slide 5: Cybersecurity Challenges"""
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Accent bar
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BOUQUET_BLUE
    accent.line.color.rgb = BOUQUET_BLUE

    # Title
    add_title_shape(slide, "The Current Cybersecurity Landscape in Mozambique")

    # Content - split into two columns
    left_content = [
        "• Critical Infrastructure Vulnerability",
        "• Government & Ministry Systems",
        "• Financial Sector Exposure",
        "• Limited Cyber Intelligence",
        "",
        "• SME & Micro-Business Risk",
        "• Cybercriminal Activity",
        "• Human Factor Weakness",
        "• Reputational Impact"
    ]

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(4.5), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    for i, item in enumerate(left_content):
        if i > 0:
            left_frame.add_paragraph()
        p = left_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    # Right side - consequences
    right_box = slide.shapes.add_shape(1, Inches(5.2), Inches(1.3), Inches(4.4), Inches(5.5))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(255, 240, 0)  # Light yellow
    right_box.line.color.rgb = GOLD
    right_box.line.width = Pt(2)

    right_text = right_box.text_frame
    right_text.word_wrap = True
    right_text.margin_top = Inches(0.2)
    right_text.margin_left = Inches(0.2)
    right_text.margin_right = Inches(0.2)

    consequences = [
        "CONSEQUENCES:",
        "",
        "• Loss of investor confidence",
        "• Increased financial losses",
        "• Compromised government security",
        "• Reduced digital competitiveness",
        "• Vulnerable critical services",
        "• Unchecked cybercrime"
    ]

    for i, item in enumerate(consequences):
        if i > 0:
            right_text.add_paragraph()
        p = right_text.paragraphs[i]
        p.text = item
        p.font.size = Pt(16)
        p.font.bold = (i == 0)
        p.font.color.rgb = BOUQUET_BLUE if i == 0 else BLACK
        p.space_before = Pt(3)
        p.space_after = Pt(3)

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text = notes_slide.notes_text_frame
    notes_text.text = """These challenges are not unique to Mozambique, but they are urgent. Without
intervention, they will worsen as digitalization increases. The good news is that
all of these challenges are addressable through a coordinated, multi-sector approach.
That's exactly what this partnership delivers. We have solutions for each of these
challenges—they just need to be deployed strategically."""

    add_footer_and_formatting(slide, 5)
    return slide

def add_awareness_matters_slide(prs, blank_layout):
    """Slide 6: Why Cybersecurity Awareness Matters"""
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Accent bar
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BOUQUET_BLUE
    accent.line.color.rgb = BOUQUET_BLUE

    # Title
    add_title_shape(slide, "The Awareness Imperative: Building a Culture of Cyber Responsibility")

    # Key stat box
    stat_box = slide.shapes.add_shape(1, Inches(0.4), Inches(1.2), Inches(9.2), Inches(0.6))
    stat_box.fill.solid()
    stat_box.fill.fore_color.rgb = RGBColor(255, 200, 0)
    stat_box.line.color.rgb = GOLD
    stat_box.line.width = Pt(3)

    stat_text = stat_box.text_frame
    stat_text.text = "85% of breaches involve human error or social engineering"
    stat_text.paragraphs[0].font.size = Pt(20)
    stat_text.paragraphs[0].font.bold = True
    stat_text.paragraphs[0].font.color.rgb = BOUQUET_BLUE
    stat_text.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Benefits
    content = [
        "Benefits of Awareness Programs:",
        "✓ Reduces vulnerability to social engineering attacks by 60-80%",
        "✓ Empowers employees to recognize and report threats",
        "✓ Creates a culture where 'security is everyone's job'",
        "✓ Protects institutional knowledge and sensitive data",
        "✓ Builds citizen resilience against cybercrime",
        "✓ Transforms cybersecurity from an IT issue to a national priority"
    ]

    add_content_box(slide, content, y=Inches(2.0), height=Inches(4.8))

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text = notes_slide.notes_text_frame
    notes_text.text = """Awareness programs are the foundation of all cybersecurity. You can have the most
sophisticated technical defenses, but if your people don't understand basic cyber
hygiene, attackers will find a way through. This program transforms every government
official, business leader, and citizen into a defender of Mozambique's cyberspace."""

    add_footer_and_formatting(slide, 6)
    return slide

def add_awareness_framework_slide(prs, blank_layout):
    """Slide 7: Awareness Program Framework"""
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Accent bar
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BOUQUET_BLUE
    accent.line.color.rgb = BOUQUET_BLUE

    # Title
    add_title_shape(slide, "National Cybersecurity Awareness Program: Framework")

    # Program tiers
    tiers = [
        ("Tier 1: Government", "Ministry staff, officials, decision-makers"),
        ("Tier 2: Critical Infrastructure", "Energy, water, communications operators"),
        ("Tier 3: Financial Sector", "Banks, payment platforms, fintech companies"),
        ("Tier 4: Private Sector & SMEs", "Large corporations, small businesses"),
        ("Tier 5: Public Awareness", "Citizens, young people, digital users")
    ]

    y_start = 1.3
    for i, (tier_name, tier_desc) in enumerate(tiers):
        y_pos = y_start + (i * 1.0)

        # Tier box
        tier_box = slide.shapes.add_shape(1, Inches(0.4), Inches(y_pos), Inches(9.2), Inches(0.9))
        tier_box.fill.solid()
        tier_box.fill.fore_color.rgb = LIGHT_GRAY
        tier_box.line.color.rgb = BOUQUET_BLUE if i % 2 == 0 else GOLD
        tier_box.line.width = Pt(2)

        text_frame = tier_box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.15)
        text_frame.margin_top = Inches(0.05)

        # Tier name
        p1 = text_frame.paragraphs[0]
        p1.text = tier_name
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = BOUQUET_BLUE

        # Tier description
        text_frame.add_paragraph()
        p2 = text_frame.paragraphs[1]
        p2.text = tier_desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = DARK_GRAY
        p2.level = 1

    # Timeline
    timeline_box = slide.shapes.add_textbox(Inches(0.4), Inches(6.3), Inches(9.2), Inches(0.8))
    timeline_frame = timeline_box.text_frame
    timeline_frame.text = "Timeline: Months 1-3 (Development) | Months 4-12 (Tier 1 & 2) | Months 13-18 (Tier 3 & 4) | Months 19-24 (Public Campaign)"
    timeline_frame.paragraphs[0].font.size = Pt(13)
    timeline_frame.paragraphs[0].font.italic = True
    timeline_frame.paragraphs[0].font.color.rgb = DARK_GRAY
    timeline_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text = notes_slide.notes_text_frame
    notes_text.text = """This is not a one-time training program—it's a sustained, multi-tiered approach to
building cyber awareness across Mozambique's entire ecosystem. We work backward from
the threats that matter most to each audience, ensuring the training is relevant,
actionable, and immediately applicable. Each tier receives specialized training
tailored to their specific risks and responsibilities."""

    add_footer_and_formatting(slide, 7)
    return slide

def add_intelligence_needs_slide(prs, blank_layout):
    """Slide 8: Cyber Intelligence Needs"""
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Accent bar
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BOUQUET_BLUE
    accent.line.color.rgb = BOUQUET_BLUE

    # Title
    add_title_shape(slide, "Cyber Intelligence Capability: Operational Needs Assessment")

    # Current gaps
    gaps_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(4.4), Inches(5.5))
    gaps_frame = gaps_box.text_frame
    gaps_frame.word_wrap = True

    gaps_title = gaps_frame.paragraphs[0]
    gaps_title.text = "Current Gaps:"
    gaps_title.font.size = Pt(18)
    gaps_title.font.bold = True
    gaps_title.font.color.rgb = BOUQUET_BLUE
    gaps_title.space_after = Pt(8)

    gaps = [
        "• No threat monitoring",
        "• Limited incident detection",
        "• No digital forensics",
        "• No early warning system",
        "• Can't track cyber criminals",
        "• No attribution capability",
        "• Reactive only (no proactive)",
        "• International gaps"
    ]

    for gap in gaps:
        gaps_frame.add_paragraph()
        p = gaps_frame.paragraphs[-1]
        p.text = gap
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(200, 0, 0)
        p.space_after = Pt(4)

    # Operational needs
    needs_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.4), Inches(5.5))
    needs_frame = needs_box.text_frame
    needs_frame.word_wrap = True

    needs_title = needs_frame.paragraphs[0]
    needs_title.text = "Needed Capabilities:"
    needs_title.font.size = Pt(18)
    needs_title.font.bold = True
    needs_title.font.color.rgb = BOUQUET_BLUE
    needs_title.space_after = Pt(8)

    needs = [
        "✓ Real-time threat monitoring",
        "✓ Rapid incident response",
        "✓ Digital forensics capability",
        "✓ Malware analysis",
        "✓ Criminal network tracking",
        "✓ Counter-terrorism support",
        "✓ Counter-organized crime",
        "✓ International cooperation"
    ]

    for need in needs:
        needs_frame.add_paragraph()
        p = needs_frame.paragraphs[-1]
        p.text = need
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0, 150, 0)
        p.space_after = Pt(4)

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text = notes_slide.notes_text_frame
    notes_text.text = """The Interior Ministry needs cyber intelligence not as a theoretical capability, but
as an operational tool that helps them protect Mozambique's most critical assets
and bring cybercriminals to justice. This requires real-time monitoring, advanced
analysis, and direct integration with investigative and law enforcement operations."""

    add_footer_and_formatting(slide, 8)
    return slide

def add_intelligence_unit_slide(prs, blank_layout):
    """Slide 9: Intelligence Unit Model"""
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Accent bar
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BOUQUET_BLUE
    accent.line.color.rgb = BOUQUET_BLUE

    # Title
    add_title_shape(slide, "Cyber Intelligence Unit: Organizational Architecture")

    # Leadership
    leadership = [
        "LEADERSHIP:",
        "• Unit Commander (Ministry-appointed)",
        "• Deputy Commander (Technical operations)",
        "• Strategic advisors (Infosec liaison)",
        "",
        "OPERATIONAL TEAMS:",
        "• Team 1: Threat Monitoring & Detection",
        "• Team 2: Digital Forensics & Analysis",
        "• Team 3: Intelligence & Attribution",
        "• Team 4: Incident Response",
        "",
        "SUPPORT:",
        "• Secure Operations Center (24/7)",
        "• Forensic Laboratory",
        "• Training Facilities",
        "• International Liaison Office"
    ]

    add_content_box(slide, leadership, y=Inches(1.2), height=Inches(5.8))

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text = notes_slide.notes_text_frame
    notes_text.text = """This unit becomes the operational core of Mozambique's cyber defense. It's not
just an intelligence service—it's an operational capability that detects threats in
real time, investigates cybercrimes, and supports law enforcement. Infosec's role
is to build this capability and then hand it off to Mozambican professionals who
will operate it independently."""

    add_footer_and_formatting(slide, 9)
    return slide

def add_cyberspace_management_slide(prs, blank_layout):
    """Slide 10: National Cyberspace Management"""
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Accent bar
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BOUQUET_BLUE
    accent.line.color.rgb = BOUQUET_BLUE

    # Title
    add_title_shape(slide, "National Cyberspace Management: Infrastructure & Governance")

    # Core Components
    components = [
        "Component 1: Centralized Threat Detection",
        "   National-level Security Operations Center (SOC)",
        "   Real-time monitoring across all sectors",
        "",
        "Component 2: Data Governance & Privacy",
        "   National data protection framework",
        "   Digital identity protection",
        "",
        "Component 3: Critical Infrastructure Protection",
        "   Monitoring of essential services",
        "   Sector-specific security baselines",
        "",
        "Component 4: Digital Sovereignty",
        "   National control over critical systems",
        "   Reduction of foreign dependencies",
        "",
        "Component 5: Incident Response Coordination",
        "   Multi-sector coordination procedures",
        "   Law enforcement integration"
    ]

    add_content_box(slide, components, y=Inches(1.2), height=Inches(5.8))

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text = notes_slide.notes_text_frame
    notes_text.text = """This architecture creates a unified national cybersecurity posture. Rather than having
each government ministry and critical infrastructure operator defending themselves
independently, Mozambique creates a coordinated defense where threats are detected
centrally, analyzed immediately, and responded to in a coordinated fashion. This is
the difference between resilience and vulnerability."""

    add_footer_and_formatting(slide, 10)
    return slide

def add_technology_slide(prs, blank_layout):
    """Slide 11: Technology & Tools"""
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Accent bar
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BOUQUET_BLUE
    accent.line.color.rgb = BOUQUET_BLUE

    # Title
    add_title_shape(slide, "Technology Stack: Tools & Platforms")

    # Technology categories
    tech_items = [
        "MONITORING & DETECTION:",
        "• SIEM Platform (Security Information & Event Management)",
        "• Network Monitoring Tools (IDS/IPS, EDR)",
        "• Vulnerability Management & Patch Management",
        "",
        "FORENSIC & ANALYSIS:",
        "• Digital Forensics Suite (disk, memory, network)",
        "• Malware Analysis Platform (sandboxed execution, behavior analysis)",
        "",
        "INTELLIGENCE & VISUALIZATION:",
        "• Threat Intelligence Platform (indicator management, attribution)",
        "• Network Mapping & Visualization (topology, relationships)",
        "",
        "SECURE COMMUNICATIONS & SUPPORT:",
        "• Encrypted Government Network",
        "• Cyber Training Platform",
        "• Knowledge Management System"
    ]

    add_content_box(slide, tech_items, y=Inches(1.2), height=Inches(5.8))

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text = notes_slide.notes_text_frame
    notes_text.text = """The technology stack is carefully selected to create an integrated, enterprise-grade
operational capability. These aren't point solutions—they work together to create a
cohesive cyber defense system. Infosec has extensive experience deploying these tools
in high-stakes environments and will ensure they're properly integrated and maintained."""

    add_footer_and_formatting(slide, 11)
    return slide

def add_roadmap_slide(prs, blank_layout):
    """Slide 12: Implementation Roadmap"""
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Accent bar
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BOUQUET_BLUE
    accent.line.color.rgb = BOUQUET_BLUE

    # Title
    add_title_shape(slide, "Implementation Timeline: 24-Month Roadmap")

    # Timeline phases
    phases = [
        ("PHASE 1: Foundation (Mo 1-3)", "Governance setup | Personnel recruitment | Program development | Infrastructure planning"),
        ("PHASE 2: Deployment (Mo 4-8)", "Awareness training begins | Cyber Intelligence Unit activated | SIEM deployed"),
        ("PHASE 3: Expansion (Mo 9-16)", "Scale awareness to all sectors | Intelligence team growth | Capability maturation"),
        ("PHASE 4: Sustainability (Mo 17-24)", "Complete training cycle | Intelligence Unit independence | National integration active")
    ]

    y_pos = 1.3
    for phase_name, phase_desc in phases:
        # Phase box
        phase_box = slide.shapes.add_shape(1, Inches(0.4), Inches(y_pos), Inches(9.2), Inches(1.2))
        phase_box.fill.solid()
        phase_box.fill.fore_color.rgb = LIGHT_GRAY
        phase_box.line.color.rgb = BOUQUET_BLUE
        phase_box.line.width = Pt(2)

        text_frame = phase_box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_right = Inches(0.2)

        p1 = text_frame.paragraphs[0]
        p1.text = phase_name
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = BOUQUET_BLUE

        text_frame.add_paragraph()
        p2 = text_frame.paragraphs[1]
        p2.text = phase_desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = DARK_GRAY

        y_pos += 1.4

    # Result box
    result_box = slide.shapes.add_textbox(Inches(0.4), Inches(6.4), Inches(9.2), Inches(0.9))
    result_frame = result_box.text_frame
    result_frame.word_wrap = True

    p = result_frame.paragraphs[0]
    p.text = "Expected Outcome: 15,000+ trained | Major cybercriminal networks disrupted | Mozambique positioned as regional cyber leader"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 150, 0)
    p.alignment = PP_ALIGN.CENTER

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text = notes_slide.notes_text_frame
    notes_text.text = """This is an aggressive but achievable timeline. The key to success is maintaining
momentum throughout the 24 months, ensuring staffing and funding remain stable, and
continuously monitoring progress against milestones. We've built in flexibility to
adjust based on real-world conditions, but the overall trajectory is clear: from
planning to fully operational national cyber capability in two years."""

    add_footer_and_formatting(slide, 12)
    return slide

def add_impact_slide(prs, blank_layout):
    """Slide 13: Expected Impact"""
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Accent bar
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BOUQUET_BLUE
    accent.line.color.rgb = BOUQUET_BLUE

    # Title
    add_title_shape(slide, "Expected Outcomes: Transforming Mozambique's Cyber Future")

    # Three outcome periods
    outcomes = [
        ("Immediate (Months 1-12)", [
            "✓ 2,000+ government officials trained",
            "✓ Real-time monitoring begins",
            "✓ First cybercrime investigations completed",
            "✓ Investor perception improves"
        ]),
        ("Medium-Term (Months 13-24)", [
            "✓ 10,000+ personnel trained",
            "✓ 70% reduction in social engineering attacks",
            "✓ Major cybercriminal networks disrupted",
            "✓ Safer digital environment established"
        ]),
        ("Long-Term (Year 2+)", [
            "✓ Cybersecurity embedded in institutions",
            "✓ National resilience demonstrated",
            "✓ Investor confidence increased",
            "✓ Regional leadership position achieved"
        ])
    ]

    y_pos = 1.3
    colors = [RGBColor(200, 220, 255), RGBColor(255, 240, 200), RGBColor(200, 255, 200)]

    for idx, (period, items) in enumerate(outcomes):
        # Period box
        period_box = slide.shapes.add_shape(1, Inches(0.4), Inches(y_pos), Inches(9.2), Inches(1.6))
        period_box.fill.solid()
        period_box.fill.fore_color.rgb = colors[idx]
        period_box.line.color.rgb = BOUQUET_BLUE
        period_box.line.width = Pt(2)

        text_frame = period_box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_right = Inches(0.2)

        p1 = text_frame.paragraphs[0]
        p1.text = period
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = BOUQUET_BLUE

        for item in items:
            text_frame.add_paragraph()
            p = text_frame.paragraphs[-1]
            p.text = item
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(2)

        y_pos += 1.8

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text = notes_slide.notes_text_frame
    notes_text.text = """This partnership transforms Mozambique from a cyber-vulnerable nation into a
cyber-resilient one. This isn't just about technology—it's about building institutional
capability, developing human expertise, and creating a culture where cybersecurity
is everyone's responsibility. The long-term impact is a nation that's better positioned
for digital economic growth and more secure against evolving threats."""

    add_footer_and_formatting(slide, 13)
    return slide

def add_partnership_rationale_slide(prs, blank_layout):
    """Slide 14: Partnership Rationale"""
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Accent bar
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BOUQUET_BLUE
    accent.line.color.rgb = BOUQUET_BLUE

    # Title
    add_title_shape(slide, "Why Infosec + Bouquet Innovation is the Right Partnership")

    # Left side: Infosec
    infosec_box = slide.shapes.add_shape(1, Inches(0.4), Inches(1.3), Inches(4.4), Inches(5.5))
    infosec_box.fill.solid()
    infosec_box.fill.fore_color.rgb = RGBColor(230, 240, 255)
    infosec_box.line.color.rgb = BOUQUET_BLUE
    infosec_box.line.width = Pt(2)

    infosec_text = infosec_box.text_frame
    infosec_text.word_wrap = True
    infosec_text.margin_left = Inches(0.15)
    infosec_text.margin_top = Inches(0.15)
    infosec_text.margin_right = Inches(0.15)

    p = infosec_text.paragraphs[0]
    p.text = "INFOSEC BRINGS:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BOUQUET_BLUE
    p.space_after = Pt(10)

    infosec_items = [
        "✓ Raw technical capability",
        "✓ Global threat intelligence",
        "✓ Forensic expertise",
        "✓ Government experience",
        "✓ 24/7 incident response",
        "✓ International partnerships",
        "✓ Advanced research labs",
        "✓ Proven track record"
    ]

    for item in infosec_items:
        infosec_text.add_paragraph()
        p = infosec_text.paragraphs[-1]
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)

    # Right side: Bouquet
    bouquet_box = slide.shapes.add_shape(1, Inches(5.2), Inches(1.3), Inches(4.4), Inches(5.5))
    bouquet_box.fill.solid()
    bouquet_box.fill.fore_color.rgb = RGBColor(255, 250, 230)
    bouquet_box.line.color.rgb = GOLD
    bouquet_box.line.width = Pt(2)

    bouquet_text = bouquet_box.text_frame
    bouquet_text.word_wrap = True
    bouquet_text.margin_left = Inches(0.15)
    bouquet_text.margin_top = Inches(0.15)
    bouquet_text.margin_right = Inches(0.15)

    p = bouquet_text.paragraphs[0]
    p.text = "BOUQUET BRINGS:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BOUQUET_BLUE
    p.space_after = Pt(10)

    bouquet_items = [
        "✓ Government relationships",
        "✓ Project management",
        "✓ Local knowledge",
        "✓ Implementation discipline",
        "✓ Mozambique expertise",
        "✓ Stakeholder coordination",
        "✓ Sustainability planning",
        "✓ Cultural bridge"
    ]

    for item in bouquet_items:
        bouquet_text.add_paragraph()
        p = bouquet_text.paragraphs[-1]
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(4)

    # Together box
    together_box = slide.shapes.add_textbox(Inches(0.4), Inches(6.9), Inches(9.2), Inches(0.6))
    together_frame = together_box.text_frame
    together_frame.text = "TOGETHER: World-class capability delivered appropriately + Sustainable, culturally-fit implementation"
    together_frame.paragraphs[0].font.size = Pt(13)
    together_frame.paragraphs[0].font.bold = True
    together_frame.paragraphs[0].font.color.rgb = RGBColor(0, 150, 0)
    together_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text = notes_slide.notes_text_frame
    notes_text.text = """This partnership isn't about choosing between technical excellence and practical delivery.
It's about having both. Infosec brings the cutting-edge capability that national security
requires. Bouquet Innovation brings the understanding of how to deploy that capability
in a Mozambican context in a way that creates lasting institutional capacity."""

    add_footer_and_formatting(slide, 14)
    return slide

def add_call_to_action_slide(prs, blank_layout):
    """Slide 15: Call to Action"""
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Accent bar
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BOUQUET_BLUE
    accent.line.color.rgb = BOUQUET_BLUE

    # Title
    add_title_shape(slide, "Moving Forward: Next Steps & Commitment")

    # Content
    content = [
        "The Opportunity:",
        "• Digital transformation is accelerating → cyber risks increase",
        "• Window of opportunity is NOW → establish resilience before threats mature",
        "• Partnership addresses FULL SPECTRUM → awareness, intelligence, infrastructure",
        "• 24-month timeline delivers QUICK RESULTS",
        "",
        "What We're Asking For:",
        "✓ Formal commitment to National Cybersecurity Council",
        "✓ Budget allocation for three-year program (~$12M total)",
        "✓ Staff designation for Cyber Intelligence Unit recruitment",
        "✓ Inter-agency coordination authorization",
        "✓ Legal/regulatory framework updates",
        "",
        "Immediate Next Steps (30 Days):",
        "→ Week 1: Government decision & authorization",
        "→ Week 2-3: Team mobilization & stakeholder meetings",
        "→ Week 4: Planning launch & recruitment begins"
    ]

    add_content_box(slide, content, y=Inches(1.2), height=Inches(5.8))

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text = notes_slide.notes_text_frame
    notes_text.text = """We're not asking for a decision today—we're asking for the conversation to continue.
If this presentation resonates with your cybersecurity challenges and strategic
objectives, we're prepared to provide much greater detail. Our teams are standing
by to answer questions, provide references, and help you understand exactly how
this partnership would work in a Mozambican context. The most important next step
is simply reaching out. We'll handle the rest."""

    add_footer_and_formatting(slide, 15)
    return slide

def add_contact_slide(prs, blank_layout):
    """Slide 16: Contact Information"""
    slide = prs.slides.add_slide(blank_layout)

    # Background with gradient effect (dark on left, light on right)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BOUQUET_BLUE

    # Accent bar
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.color.rgb = GOLD

    # Title (white on blue background)
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.4), Inches(9.2), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Contact Information & Partnership Inquiry"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Bouquet section
    bouquet_section = slide.shapes.add_shape(1, Inches(0.4), Inches(1.4), Inches(4.4), Inches(5.2))
    bouquet_section.fill.solid()
    bouquet_section.fill.fore_color.rgb = WHITE
    bouquet_section.line.color.rgb = BOUQUET_BLUE
    bouquet_section.line.width = Pt(3)

    bouquet_text = bouquet_section.text_frame
    bouquet_text.word_wrap = True
    bouquet_text.margin_left = Inches(0.2)
    bouquet_text.margin_top = Inches(0.2)
    bouquet_text.margin_right = Inches(0.2)

    p1 = bouquet_text.paragraphs[0]
    p1.text = "Bouquet Innovation"
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = BOUQUET_BLUE
    p1.space_after = Pt(10)

    bouquet_items = [
        "Email:",
        "support@bouquet-innovation.net",
        "",
        "Phone:",
        "+258 878 275 656",
        "",
        "Focus:",
        "Digital transformation,",
        "Government modernization,",
        "Cybersecurity strategy",
        "",
        "Expertise:",
        "Portuguese-speaking Africa,",
        "Government partnerships"
    ]

    for item in bouquet_items:
        bouquet_text.add_paragraph()
        p = bouquet_text.paragraphs[-1]
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY if item and item[0] not in "ES" else BOUQUET_BLUE
        if item in ["Email:", "Phone:", "Focus:", "Expertise:"]:
            p.font.bold = True
        p.space_after = Pt(2)

    # Infosec section
    infosec_section = slide.shapes.add_shape(1, Inches(5.2), Inches(1.4), Inches(4.4), Inches(5.2))
    infosec_section.fill.solid()
    infosec_section.fill.fore_color.rgb = WHITE
    infosec_section.line.color.rgb = GOLD
    infosec_section.line.width = Pt(3)

    infosec_text = infosec_section.text_frame
    infosec_text.word_wrap = True
    infosec_text.margin_left = Inches(0.2)
    infosec_text.margin_top = Inches(0.2)
    infosec_text.margin_right = Inches(0.2)

    p1 = infosec_text.paragraphs[0]
    p1.text = "Infosec"
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = BOUQUET_BLUE
    p1.space_after = Pt(10)

    infosec_items = [
        "Global Cyber Threat Leader",
        "",
        "Specializations:",
        "• Cyber intelligence",
        "• Digital forensics",
        "• Incident response",
        "• Threat research",
        "",
        "Government Experience:",
        "• Trusted by national",
        "  security agencies",
        "• 24/7 incident response",
        "• Six continent network"
    ]

    for item in infosec_items:
        infosec_text.add_paragraph()
        p = infosec_text.paragraphs[-1]
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(2)

    # Call to action box at bottom
    cta_box = slide.shapes.add_textbox(Inches(0.4), Inches(6.9), Inches(9.2), Inches(0.6))
    cta_frame = cta_box.text_frame
    cta_frame.text = "Email: support@bouquet-innovation.net | Phone: +258 878 275 656 | Ready to discuss partnership opportunities"
    cta_frame.paragraphs[0].font.size = Pt(14)
    cta_frame.paragraphs[0].font.bold = True
    cta_frame.paragraphs[0].font.color.rgb = GOLD
    cta_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_text = notes_slide.notes_text_frame
    notes_text.text = """We're not asking for a decision today—we're asking for the conversation to continue.
If this presentation resonates with your cybersecurity challenges and strategic
objectives, we're prepared to provide much greater detail. Our teams are standing
by to answer questions, provide references, and help you understand exactly how
this partnership would work in a Mozambican context. Contact us to schedule a
detailed briefing with both Bouquet Innovation and Infosec leadership."""

    add_footer_and_formatting(slide, 16)
    return slide

if __name__ == "__main__":
    create_presentation()
